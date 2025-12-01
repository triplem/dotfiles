"""Service for converting linked files to markdown using lightweight libraries."""
import logging
import hashlib
from pathlib import Path
from typing import Optional, Dict, Any

logger = logging.getLogger(__name__)


class DocumentConverter:
    """Service for converting documents to markdown using lightweight libraries.

    Priority handlers (faster, tested):
    - PDF files via pymupdf4llm
    - DOCX files via python-docx
    - PPTX files via python-pptx

    Fallback handler via markitdown for:
    - Excel (.xlsx, .xls, .csv)
    - Images (.png, .jpg, .jpeg, .gif, .bmp)
    - Audio (.mp3, .wav, .m4a)
    - HTML (.html, .htm)
    - Text formats (.json, .xml)
    - ZIP files (.zip)
    - And more...
    """

    # Priority formats (use specialized libraries)
    PRIORITY_EXTENSIONS = {
        '.pdf',   # pymupdf4llm (fastest for PDFs)
        '.docx',  # python-docx (fastest for Word)
        '.pptx',  # python-pptx (fastest for PowerPoint)
    }

    # Audio formats (can cause multiprocessing leaks in markitdown)
    AUDIO_EXTENSIONS = {'.mp3', '.wav', '.m4a', '.ogg', '.flac'}

    # Supported via markitdown (fallback)
    MARKITDOWN_EXTENSIONS = {
        '.xlsx', '.xls', '.csv',  # Excel/spreadsheets
        '.html', '.htm',          # Web pages
        '.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff',  # Images
        '.json', '.xml',          # Text formats
        '.zip',                   # Archives
        '.epub',                  # eBooks
        '.ppt',                   # Legacy PowerPoint (via markitdown)
        '.doc',                   # Legacy Word (via markitdown)
    } | AUDIO_EXTENSIONS  # Include audio in markitdown extensions

    # All supported extensions
    SUPPORTED_EXTENSIONS = PRIORITY_EXTENSIONS | MARKITDOWN_EXTENSIONS

    def __init__(self, skip_audio: bool = True):
        """Initialize the document conversion service.

        Args:
            skip_audio: If True, skip audio files to avoid multiprocessing leaks from markitdown.
        """
        self.skip_audio = skip_audio
        logger.info(f"DocumentConverter initialized (priority: pymupdf4llm/python-docx/python-pptx, fallback: markitdown, skip_audio: {skip_audio})")

    @staticmethod
    def calculate_md5(file_path: str) -> str:
        """Calculate MD5 hash of a file for change detection."""
        hash_md5 = hashlib.md5()
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(4096), b""):
                hash_md5.update(chunk)
        return hash_md5.hexdigest()

    def is_supported(self, file_path: str) -> bool:
        """Check if a file extension is supported."""
        ext = Path(file_path).suffix.lower()
        is_supported = ext in self.SUPPORTED_EXTENSIONS
        if not is_supported and ext:
            logger.warning(f"File type '{ext}' is not supported")
        return is_supported

    def convert_to_markdown(
        self,
        file_path: str,
        max_file_size: int = 52428800  # 50MB default
    ) -> Dict[str, Any]:
        """
        Convert a document to markdown using specialized libraries.

        Args:
            file_path: Path to the file to convert
            max_file_size: Maximum file size in bytes (default 50MB)

        Returns:
            Dictionary with:
                - status: 'success', 'error', 'skipped'
                - markdown: The converted markdown text (if successful)
                - md5: MD5 hash of the source file
                - error: Error message (if failed)
                - file_size: Size of the file in bytes
        """
        path = Path(file_path)

        # Check file exists
        if not path.exists():
            logger.warning(f"File not found: {file_path}")
            return {
                "status": "error",
                "error": "File not found",
                "md5": None,
                "file_size": 0
            }

        # Check file size
        file_size = path.stat().st_size
        if file_size > max_file_size:
            logger.warning(f"File too large: {file_path} ({file_size} bytes)")
            return {
                "status": "skipped",
                "error": f"File too large: {file_size / 1024 / 1024:.1f}MB",
                "md5": None,
                "file_size": file_size
            }

        # Check extension
        ext = path.suffix.lower()
        if not self.is_supported(file_path):
            logger.warning(f"Unsupported file extension: {ext}")
            return {
                "status": "error",
                "error": f"Unsupported file extension: {ext}",
                "md5": None,
                "file_size": file_size
            }

        # Skip audio files if configured (to avoid multiprocessing leaks)
        if self.skip_audio and ext in self.AUDIO_EXTENSIONS:
            logger.info(f"Skipping audio file {file_path} (skip_audio=True)")
            return {
                "status": "skipped",
                "error": "Audio files skipped to avoid multiprocessing leaks",
                "md5": None,
                "file_size": file_size
            }

        # Calculate MD5
        try:
            md5 = self.calculate_md5(file_path)
        except Exception as e:
            logger.error(f"Error calculating MD5 for {file_path}: {e}")
            return {
                "status": "error",
                "error": f"Error calculating MD5: {str(e)}",
                "md5": None,
                "file_size": file_size
            }

        # Convert document based on extension
        ext = path.suffix.lower()

        try:
            # Use priority handlers for known formats (faster)
            if ext == '.pdf':
                return self._convert_pdf_with_pymupdf(file_path, md5, file_size)
            elif ext == '.docx':
                return self._convert_docx_with_python_docx(file_path, md5, file_size)
            elif ext == '.pptx':
                return self._convert_pptx_with_python_pptx(file_path, md5, file_size)
            elif ext in self.MARKITDOWN_EXTENSIONS:
                # Use markitdown for other supported formats
                return self._convert_with_markitdown(file_path, md5, file_size)
            else:
                # Unsupported format
                logger.warning(f"Cannot convert {file_path}: unsupported file type '{ext}'")
                return {
                    "status": "error",
                    "error": f"Unsupported file type: {ext}",
                    "md5": md5,
                    "file_size": file_size
                }
        except Exception as e:
            logger.error(f"Unexpected error converting {file_path}: {e}", exc_info=True)
            return {
                "status": "error",
                "error": f"Conversion failed: {str(e)}",
                "md5": md5,
                "file_size": file_size
            }

    def _convert_pdf_with_pymupdf(
        self,
        file_path: str,
        md5: str,
        file_size: int
    ) -> Dict[str, Any]:
        """Convert PDF using lightweight pymupdf4llm (fast, low memory)."""
        try:
            import pymupdf4llm

            logger.info(f"Converting PDF {file_path} with pymupdf4llm...")
            markdown_text = pymupdf4llm.to_markdown(file_path)

            logger.info(f"Successfully converted {file_path} ({len(markdown_text)} chars)")
            return {
                "status": "success",
                "markdown": markdown_text,
                "md5": md5,
                "file_size": file_size
            }

        except Exception as e:
            logger.error(f"Error converting PDF {file_path}: {e}", exc_info=True)
            return {
                "status": "error",
                "error": str(e),
                "md5": md5,
                "file_size": file_size
            }

    def _convert_docx_with_python_docx(
        self,
        file_path: str,
        md5: str,
        file_size: int
    ) -> Dict[str, Any]:
        """Convert DOCX using lightweight python-docx (fast, low memory)."""
        try:
            from docx import Document

            logger.info(f"Converting DOCX {file_path} with python-docx...")
            doc = Document(file_path)
            text = "\n".join([para.text for para in doc.paragraphs])

            logger.info(f"Successfully converted {file_path} ({len(text)} chars)")
            return {
                "status": "success",
                "markdown": text,  # Plain text, not markdown, but good enough
                "md5": md5,
                "file_size": file_size
            }

        except Exception as e:
            logger.error(f"Error converting DOCX {file_path}: {e}", exc_info=True)
            return {
                "status": "error",
                "error": str(e),
                "md5": md5,
                "file_size": file_size
            }

    def _convert_pptx_with_python_pptx(
        self,
        file_path: str,
        md5: str,
        file_size: int
    ) -> Dict[str, Any]:
        """Convert PPTX using lightweight python-pptx (fast, low memory)."""
        try:
            from pptx import Presentation

            logger.info(f"Converting PPTX {file_path} with python-pptx...")
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)

            result_text = "\n".join(text)

            logger.info(f"Successfully converted {file_path} ({len(result_text)} chars)")
            return {
                "status": "success",
                "markdown": result_text,  # Plain text, not markdown, but good enough
                "md5": md5,
                "file_size": file_size
            }

        except Exception as e:
            logger.error(f"Error converting PPTX {file_path}: {e}", exc_info=True)
            return {
                "status": "error",
                "error": str(e),
                "md5": md5,
                "file_size": file_size
            }

    def _convert_with_markitdown(
        self,
        file_path: str,
        md5: str,
        file_size: int
    ) -> Dict[str, Any]:
        """Convert file using markitdown (fallback for various formats).

        Creates a fresh MarkItDown instance for each conversion to avoid
        memory leaks from multiprocessing (ASR for audio files).
        """
        try:
            from markitdown import MarkItDown
            import gc

            ext = Path(file_path).suffix.lower()
            logger.info(f"Converting {ext} file {file_path} with markitdown...")

            # Create fresh instance for each conversion to avoid resource leaks
            md_converter = MarkItDown()

            try:
                result = md_converter.convert(file_path)
                markdown_text = result.text_content

                logger.info(f"Successfully converted {file_path} ({len(markdown_text)} chars)")
                return {
                    "status": "success",
                    "markdown": markdown_text,
                    "md5": md5,
                    "file_size": file_size
                }
            finally:
                # Clean up to prevent multiprocessing resource leaks
                del md_converter
                gc.collect()

        except Exception as e:
            logger.error(f"Error converting {file_path} with markitdown: {e}", exc_info=True)
            return {
                "status": "error",
                "error": str(e),
                "md5": md5,
                "file_size": file_size
            }


# Global instance
_document_converter: Optional[DocumentConverter] = None


def get_document_converter() -> DocumentConverter:
    """Get or create the global document conversion service instance."""
    global _document_converter
    if _document_converter is None:
        _document_converter = DocumentConverter(skip_audio=True)
    return _document_converter
